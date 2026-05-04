"""
Generate Solaroo RE CRM flowcharts as a .pptx deck.

Slides:
  1. Title
  2. Master end-to-end lifecycle (1-page overview, swim-lane style)
  3. Lead capture & qualification
  4. Opportunity stage flow with discipline layer (next-action / health / activities)
  5. Proposal lifecycle (versions, assumptions, approval)
  6. Award & handover (Opportunity WON -> Contract -> Project)
  7. Project execution with gates + RAG/blocker flow
  8. Procurement flow (RFQ -> PO -> Delivery)
  9. Commissioning -> Asset Register -> O&M handover
  10. Cross-cutting: notifications, audit, role visibility

Uses python-pptx shapes + connectors.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUT = os.path.join(os.path.dirname(__file__), "Solaroo_CRM_Flowcharts.pptx")

# ─── Palette (Monday-ish) ─────────────────────────────────────────────────────
C_BLUE     = RGBColor(0x00, 0x73, 0xEA)
C_BLUE_LT  = RGBColor(0xDC, 0xE9, 0xFC)
C_GREEN    = RGBColor(0x00, 0xCA, 0x72)
C_GREEN_LT = RGBColor(0xD7, 0xF3, 0xE3)
C_RED      = RGBColor(0xE2, 0x44, 0x5C)
C_RED_LT   = RGBColor(0xFB, 0xDC, 0xE0)
C_AMBER    = RGBColor(0xFD, 0xAB, 0x3D)
C_AMBER_LT = RGBColor(0xFE, 0xEA, 0xC8)
C_PURPLE   = RGBColor(0xA2, 0x5D, 0xDC)
C_PURPLE_LT= RGBColor(0xEA, 0xDB, 0xF8)
C_TEAL     = RGBColor(0x00, 0x86, 0x9B)
C_TEAL_LT  = RGBColor(0xCC, 0xE7, 0xEB)
C_GREY     = RGBColor(0x67, 0x68, 0x79)
C_GREY_LT  = RGBColor(0xF0, 0xF1, 0xF4)
C_DARK     = RGBColor(0x32, 0x33, 0x38)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9 widescreen
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Monkey-patch shape/connector/textbox factories to coerce floats to int ──
# OOXML EMU coordinates must be integers; division (`/`) in calculations
# produces floats which python-pptx serializes as "1.5" — invalid OOXML.
from pptx.shapes.shapetree import SlideShapes as _SS

def _coerce_args(args, n_to_coerce):
    args = list(args)
    for i in range(min(n_to_coerce, len(args))):
        args[i] = int(round(float(args[i])))
    return tuple(args)

_orig_add_shape    = _SS.add_shape
_orig_add_textbox  = _SS.add_textbox
_orig_add_connector = _SS.add_connector

def _patched_add_shape(self, autoshape_type_id, x, y, cx, cy):
    return _orig_add_shape(self, autoshape_type_id,
                           int(round(float(x))), int(round(float(y))),
                           int(round(float(cx))), int(round(float(cy))))

def _patched_add_textbox(self, x, y, w, h):
    return _orig_add_textbox(self,
                             int(round(float(x))), int(round(float(y))),
                             int(round(float(w))), int(round(float(h))))

def _patched_add_connector(self, connector_type, begin_x, begin_y, end_x, end_y):
    return _orig_add_connector(self, connector_type,
                               int(round(float(begin_x))), int(round(float(begin_y))),
                               int(round(float(end_x))),   int(round(float(end_y))))

_SS.add_shape    = _patched_add_shape
_SS.add_textbox  = _patched_add_textbox
_SS.add_connector = _patched_add_connector

BLANK = prs.slide_layouts[6]

# ─── Helpers ──────────────────────────────────────────────────────────────────

def _i(v):
    """Coerce EMU value to int — OOXML requires integer coordinates."""
    return int(round(float(v)))

def add_text(slide, x, y, w, h, text, *, size=11, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb

def add_box(slide, x, y, w, h, text, *, fill=C_BLUE, border=None, text_color=C_WHITE,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, size=11, bold=True):
    x, y, w, h = _i(x), _i(y), _i(w), _i(h)
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border if border else fill
    s.line.width = Pt(0.75)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    r.font.name = "Calibri"
    return s

def add_diamond(slide, x, y, w, h, text, *, fill=C_AMBER, text_color=C_WHITE, size=10):
    return add_box(slide, x, y, w, h, text, fill=fill, text_color=text_color,
                   shape=MSO_SHAPE.DIAMOND, size=size, bold=True)

from pptx.oxml.ns import qn
from lxml import etree

def _set_tail_arrow(line, color, weight):
    """Set line color/width AND a tail arrowhead, ensuring schema order."""
    ln = line._get_or_add_ln()
    # Remove any auto-added fill / arrow children so we control order
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill',
                'a:prstDash', 'a:custDash',
                'a:round', 'a:bevel', 'a:miter',
                'a:headEnd', 'a:tailEnd'):
        for el in ln.findall(qn(tag)):
            ln.remove(el)
    # Set width on <a:ln>
    ln.set('w', str(int(Pt(weight))))
    ln.set('cap', 'flat')
    ln.set('cmpd', 'sng')
    ln.set('algn', 'ctr')
    # Add children in schema order: fill, dash, join, headEnd, tailEnd
    fill = etree.SubElement(ln, qn('a:solidFill'))
    clr = etree.SubElement(fill, qn('a:srgbClr'))
    clr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))
    dash = etree.SubElement(ln, qn('a:prstDash'))
    dash.set('val', 'solid')
    etree.SubElement(ln, qn('a:round'))
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')

def add_arrow(slide, x1, y1, x2, y2, *, color=C_GREY, weight=1.75):
    x1, y1, x2, y2 = _i(x1), _i(y1), _i(x2), _i(y2)
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    _set_tail_arrow(conn.line, color, weight)
    return conn

def add_elbow_arrow(slide, x1, y1, x2, y2, *, color=C_GREY, weight=1.75):
    x1, y1, x2, y2 = _i(x1), _i(y1), _i(x2), _i(y2)
    conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, x1, y1, x2, y2)
    _set_tail_arrow(conn.line, color, weight)
    return conn

def slide_title(slide, title, subtitle=None):
    # Top color band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.55))
    band.fill.solid(); band.fill.fore_color.rgb = C_BLUE
    band.line.fill.background()
    band.shadow.inherit = False
    add_text(slide, Inches(0.4), Inches(0.05), Inches(12), Inches(0.45),
             title, size=18, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.6), Inches(12), Inches(0.35),
                 subtitle, size=11, color=C_GREY)
    # Footer
    add_text(slide, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.3),
             "Solaroo RE CRM — Sales-to-Project Lifecycle Flowcharts", size=8, color=C_GREY)

def add_legend(slide, x, y, items):
    """items = [(label, color), ...]"""
    cx = x
    for label, color in items:
        sw = Inches(0.18)
        sh = Inches(0.18)
        sw_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y, sw, sh)
        sw_box.fill.solid(); sw_box.fill.fore_color.rgb = color
        sw_box.line.color.rgb = color
        add_text(slide, cx + Inches(0.22), y - Inches(0.02), Inches(1.6), Inches(0.22),
                 label, size=9, color=C_DARK)
        cx += Inches(1.7)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
# Background gradient strip
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
band.fill.solid(); band.fill.fore_color.rgb = C_GREY_LT
band.line.fill.background()
band.shadow.inherit = False

side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), SLIDE_H)
side.fill.solid(); side.fill.fore_color.rgb = C_BLUE
side.line.fill.background()
side.shadow.inherit = False

add_text(s, Inches(0.9), Inches(2.4), Inches(11), Inches(1),
         "Solaroo RE CRM", size=44, bold=True, color=C_DARK)
add_text(s, Inches(0.9), Inches(3.3), Inches(11), Inches(0.6),
         "Sales-to-Project Lifecycle Flowcharts", size=24, color=C_BLUE)
add_text(s, Inches(0.9), Inches(4.1), Inches(11), Inches(0.4),
         "Account · Site · Opportunity · Proposal · Contract · Project · Commissioning · O&M",
         size=13, color=C_GREY)

add_text(s, Inches(0.9), Inches(6.6), Inches(11), Inches(0.3),
         "Pekat Teknologi Sdn Bhd · Internal CRM platform · v1", size=10, color=C_GREY)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 2 — Master end-to-end flow (horizontal swim-lane summary)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "End-to-End Lifecycle — One-Page Master Flow",
            "From first contact to operating asset. Each stage is a CRM module.")

# 9 stage boxes across the slide
stages = [
    ("Account\n& Site",       C_TEAL,   "Customer + location\nrecord created"),
    ("Contact",               C_TEAL,   "Decision-makers\nlinked to account"),
    ("Opportunity",           C_BLUE,   "Deal record +\nstage history"),
    ("Proposal\n(versions)",  C_BLUE,   "Frozen assumption\nset per version"),
    ("Contract\nawarded",     C_GREEN,  "Signed scope +\nmilestones"),
    ("Project\n(auto-created)", C_GREEN, "PM assigned at WON.\nGates open."),
    ("Procurement\n+ Site Exec", C_PURPLE, "RFQs · POs ·\nDeliveries · Build"),
    ("Commissioning\n+ Handover", C_AMBER, "Tests passed.\nAssets registered."),
    ("O&M\n(operating)",      C_AMBER,  "Tickets, maintenance,\nwarranty"),
]

n = len(stages)
gap = Inches(0.1)
total_w = SLIDE_W - Inches(0.8)
box_w = (total_w - gap * (n - 1)) / n
box_h = Inches(1.05)
y_top = Inches(1.6)

for i, (label, color, sub) in enumerate(stages):
    x = Inches(0.4) + (box_w + gap) * i
    add_box(s, x, y_top, box_w, box_h, label, fill=color, size=11)
    add_text(s, x, y_top + box_h + Inches(0.05), box_w, Inches(0.7),
             sub, size=8, color=C_GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# Connecting arrows beneath the boxes
arrow_y = y_top + box_h / 2
for i in range(n - 1):
    x1 = Inches(0.4) + (box_w + gap) * i + box_w
    x2 = Inches(0.4) + (box_w + gap) * (i + 1)
    add_arrow(s, x1, arrow_y, x2, arrow_y, color=C_GREY, weight=1.5)

# Discipline / cross-cutting band
band_y = Inches(4.2)
add_text(s, Inches(0.4), band_y, Inches(12.5), Inches(0.3),
         "Cross-cutting — applied to every stage:", size=11, bold=True, color=C_DARK)

cross = [
    ("Activity Timeline",      C_BLUE_LT,  C_BLUE,
     "Every call/email/visit logged"),
    ("Next-Action Discipline", C_AMBER_LT, C_AMBER,
     "One open action per opp"),
    ("Health Indicator",       C_GREEN_LT, C_GREEN,
     "Healthy / At Risk / Stale / Overdue"),
    ("Notifications",          C_PURPLE_LT,C_PURPLE,
     "Bell alerts + daily digest"),
    ("Audit Log",              C_GREY_LT,  C_GREY,
     "Every status change recorded"),
    ("Role-based Visibility",  C_RED_LT,   C_RED,
     "Margin/CAPEX restricted"),
]
cw = (total_w - gap * (len(cross) - 1)) / len(cross)
for i, (label, fill, border, sub) in enumerate(cross):
    x = Inches(0.4) + (cw + gap) * i
    box = add_box(s, x, Inches(4.55), cw, Inches(0.55), label,
                  fill=fill, text_color=border, border=border, size=10)
    add_text(s, x, Inches(5.15), cw, Inches(0.5), sub, size=8,
             color=C_GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# Roles band
add_text(s, Inches(0.4), Inches(5.95), Inches(12.5), Inches(0.3),
         "Primary owners by stage:", size=11, bold=True, color=C_DARK)

owners = [
    "Sales Engineer", "Sales Engineer", "Sales Engineer / Manager",
    "Sales Eng + Design Lead", "Sales Manager + Director",
    "PMO Manager → PM", "Procurement + Site Sup.",
    "Commissioning Eng.", "O&M Engineer",
]
for i, owner in enumerate(owners):
    x = Inches(0.4) + (box_w + gap) * i
    add_text(s, x, Inches(6.3), box_w, Inches(0.5),
             owner, size=8, color=C_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 3 — Lead capture & qualification
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "1. Lead Capture & Qualification",
            "How a new enquiry becomes a qualified opportunity in the CRM.")

# Top: inputs (lead sources)
add_text(s, Inches(0.4), Inches(1.2), Inches(3), Inches(0.3),
         "Lead sources", size=11, bold=True, color=C_GREY)

sources = ["Inbound\nweb / call", "Referral", "Trade show /\nexpo", "Existing\naccount"]
for i, src in enumerate(sources):
    add_box(s, Inches(0.4 + i*1.55), Inches(1.55), Inches(1.4), Inches(0.7),
            src, fill=C_TEAL_LT, text_color=C_TEAL, border=C_TEAL, size=9)

# Funnel into "Capture in CRM"
capture = add_box(s, Inches(0.4), Inches(2.7), Inches(7), Inches(0.55),
                  "Sales Engineer captures lead in CRM",
                  fill=C_BLUE, size=11)

for i in range(4):
    add_arrow(s, Inches(1.1 + i*1.55), Inches(2.25), Inches(1.1 + i*1.55), Inches(2.7))

# Three parallel records created
y_rec = Inches(3.6)
recs = [
    ("Account\n(if new)",   "company / industry"),
    ("Site\n(if new)",      "geo, capacity, grid"),
    ("Contact(s)",          "name, role, email"),
]
rw = Inches(2.1); rgap = Inches(0.25)
for i, (label, sub) in enumerate(recs):
    x = Inches(0.4) + (rw + rgap) * i
    add_box(s, x, y_rec, rw, Inches(0.55), label, fill=C_TEAL, size=10)
    add_text(s, x, y_rec + Inches(0.6), rw, Inches(0.4), sub,
             size=8, color=C_GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    add_arrow(s, Inches(3.9), Inches(3.25), x + rw/2, y_rec)

# Opportunity created
opp_y = Inches(4.85)
add_box(s, Inches(0.4), opp_y, Inches(7), Inches(0.55),
        "Opportunity created · stage = QUALIFICATION",
        fill=C_BLUE, size=11)
for i in range(3):
    x = Inches(0.4) + (rw + rgap) * i + rw/2
    add_arrow(s, x, y_rec + Inches(1.0), x, opp_y)

# Decision diamond: qualified?
dia = add_diamond(s, Inches(2.3), Inches(5.7), Inches(3.2), Inches(1.1),
                  "Qualified?\n(budget · authority ·\nneed · timeline)",
                  fill=C_AMBER, size=10)
add_arrow(s, Inches(3.9), opp_y + Inches(0.55), Inches(3.9), Inches(5.7))

# YES → move to needs assessment
add_box(s, Inches(6.0), Inches(5.85), Inches(2.6), Inches(0.6),
        "Stage → NEEDS_ASSESSMENT", fill=C_GREEN, size=10)
add_arrow(s, Inches(5.5), Inches(6.25), Inches(6.0), Inches(6.25), color=C_GREEN, weight=2)
add_text(s, Inches(5.55), Inches(6.0), Inches(0.45), Inches(0.25), "Yes",
         size=9, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# NO → mark LOST
add_box(s, Inches(2.3), Inches(7.0), Inches(3.2), Inches(0.4),
        "Stage → LOST  ·  reason recorded", fill=C_RED, size=9)
add_arrow(s, Inches(3.9), Inches(6.8), Inches(3.9), Inches(7.0), color=C_RED, weight=2)
add_text(s, Inches(4.0), Inches(6.8), Inches(0.4), Inches(0.25), "No",
         size=9, bold=True, color=C_RED, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# Right-side discipline panel
panel_x = Inches(9.3); panel_y = Inches(1.2)
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           panel_x, panel_y, Inches(3.6), Inches(5.8))
panel.fill.solid(); panel.fill.fore_color.rgb = C_GREY_LT
panel.line.color.rgb = C_GREY
panel.shadow.inherit = False
add_text(s, panel_x + Inches(0.2), panel_y + Inches(0.15),
         Inches(3.2), Inches(0.3),
         "Discipline applied immediately",
         size=12, bold=True, color=C_DARK)

bullets = [
    "• Set Owner = capturing user",
    "• Log inbound activity (CALL / EMAIL)",
    "• Set first Next Action: SITE_SURVEY or FOLLOW_UP",
    "• Set Due Date (≤ 7 days)",
    "• Save → opportunity becomes 'Healthy'",
    "",
    "Notifications wired:",
    "• Overdue next action (daily)",
    "• No activity in 14 / 30 days",
    "• Proposal awaiting follow-up (3d)",
]
y = panel_y + Inches(0.55)
for line in bullets:
    add_text(s, panel_x + Inches(0.25), y, Inches(3.2), Inches(0.3),
             line, size=10, color=C_DARK, anchor=MSO_ANCHOR.TOP)
    y += Inches(0.36)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 4 — Opportunity stage flow with discipline
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "2. Opportunity Stage Flow",
            "Linear stage progression. Stage jumps are blocked. Every transition is audited.")

# 8 stages along a centred horizontal axis
stages = [
    "QUALIFICATION", "NEEDS_ASSESSMENT", "TECHNICAL_DESIGN",
    "BUDGETARY_PROPOSAL", "FIRM_PROPOSAL", "NEGOTIATION",
    "CLOSING", "WON",
]
n = len(stages)
gap = Inches(0.12)
total_w = SLIDE_W - Inches(0.8)
sw = (total_w - gap * (n - 1)) / n
sh = Inches(0.7)
y_axis = Inches(2.0)

for i, label in enumerate(stages):
    x = Inches(0.4) + (sw + gap) * i
    color = C_GREEN if label == "WON" else C_BLUE
    add_box(s, x, y_axis, sw, sh, label, fill=color, size=9)
    if i < n - 1:
        add_arrow(s, x + sw, y_axis + sh/2,
                  x + sw + gap, y_axis + sh/2)

# LOST exit branch (off-axis)
lost_y = Inches(3.4)
add_box(s, Inches(5.5), lost_y, Inches(2.3), Inches(0.55),
        "LOST · reason required", fill=C_RED, size=10)
# arrow from middle of axis down to LOST
add_elbow_arrow(s, Inches(6.65), y_axis + sh, Inches(6.65), lost_y, color=C_RED, weight=1.75)
add_text(s, Inches(6.7), Inches(2.85), Inches(2), Inches(0.3),
         "any stage → LOST", size=9, color=C_RED, anchor=MSO_ANCHOR.MIDDLE)

# Per-stage discipline lane below
disc_y = Inches(4.3)
add_text(s, Inches(0.4), Inches(4.0), Inches(12.5), Inches(0.3),
         "What the system enforces at every stage:", size=11, bold=True, color=C_DARK)

discs = [
    ("Activity logging",   C_BLUE_LT,  C_BLUE,
     "Sales engineer logs call / email / visit / WhatsApp / proposal follow-up"),
    ("Next action set",    C_AMBER_LT, C_AMBER,
     "Type · owner · due date.  Empty slot ⇒ At Risk flag"),
    ("Health calc",        C_GREEN_LT, C_GREEN,
     "Healthy · At Risk (>14d silent / no action) · Stale (>30d) · Overdue (action past due)"),
    ("Notifications",      C_PURPLE_LT,C_PURPLE,
     "Bell alerts to owner + sales manager.  Daily digest at 30-min tick"),
    ("Audit",              C_GREY_LT,  C_GREY,
     "Every stage change, every field touch, recorded in audit log"),
]
dh = Inches(0.55); dgap = Inches(0.1)
for i, (label, fill, border, sub) in enumerate(discs):
    y = disc_y + (dh + dgap) * i
    add_box(s, Inches(0.4), y, Inches(2.5), dh, label,
            fill=fill, text_color=border, border=border, size=10)
    add_text(s, Inches(3.0), y, Inches(10), dh, sub, size=10, color=C_DARK)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 5 — Proposal lifecycle
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "3. Proposal Lifecycle (versions & approval)",
            "Each proposal version freezes its assumption set. Approved versions are immutable.")

# Top swim lane: roles
lanes = [
    ("Sales Engineer",    C_BLUE,    Inches(1.1)),
    ("Design Lead",       C_TEAL,    Inches(2.5)),
    ("Director / Sales Mgr (Approver)", C_PURPLE, Inches(3.9)),
    ("Customer",          C_GREEN,   Inches(5.3)),
]
for label, color, y in lanes:
    # lane background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), y,
                            Inches(11.0), Inches(1.3))
    bg.fill.solid(); bg.fill.fore_color.rgb = C_GREY_LT
    bg.line.color.rgb = C_WHITE
    bg.shadow.inherit = False
    # lane label
    lbl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), y,
                             Inches(1.55), Inches(1.3))
    lbl.fill.solid(); lbl.fill.fore_color.rgb = color
    lbl.line.color.rgb = color
    lbl.shadow.inherit = False
    tf = lbl.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = C_WHITE
    r.font.name = "Calibri"

# Sales Engineer steps
add_box(s, Inches(2.2), Inches(1.35), Inches(1.7), Inches(0.7),
        "Draft v1\n+ assumptions", fill=C_BLUE, size=9)
add_box(s, Inches(7.7), Inches(1.35), Inches(1.7), Inches(0.7),
        "Revise v2 if\nrejected", fill=C_BLUE, size=9)
add_box(s, Inches(11.0), Inches(1.35), Inches(1.7), Inches(0.7),
        "Issue v_n\nto customer", fill=C_BLUE, size=9)

# Design Lead
add_box(s, Inches(4.0), Inches(2.75), Inches(1.7), Inches(0.7),
        "Tech sign-off\non BoQ / SLD", fill=C_TEAL, size=9)

# Approver
add_box(s, Inches(5.8), Inches(4.15), Inches(1.7), Inches(0.7),
        "Approve\nor reject", fill=C_PURPLE, size=9)

# Customer
add_box(s, Inches(11.0), Inches(5.55), Inches(1.7), Inches(0.7),
        "Accept /\nrequest revision", fill=C_GREEN, size=9)

# Arrows
add_arrow(s, Inches(3.05), Inches(2.05), Inches(4.85), Inches(2.75))
add_arrow(s, Inches(4.85), Inches(3.45), Inches(6.65), Inches(4.15))
add_arrow(s, Inches(6.65), Inches(4.85), Inches(11.85), Inches(5.55), color=C_GREEN)
# Reject loop back
add_arrow(s, Inches(5.8), Inches(4.5), Inches(3.9), Inches(2.05), color=C_RED)
add_text(s, Inches(4.4), Inches(3.75), Inches(1.5), Inches(0.3),
         "reject → new version", size=9, color=C_RED)
# Customer revision back
add_arrow(s, Inches(11.0), Inches(5.7), Inches(8.55), Inches(2.05), color=C_AMBER)
add_text(s, Inches(9.0), Inches(3.7), Inches(2.0), Inches(0.3),
         "revision requested",
         size=9, color=C_AMBER)

# Bottom: rules
add_text(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3),
         "Rules: each version freezes its assumption set · approved versions are immutable · "
         "supersession requires a new version · all decisions audited · margin/CAPEX visible to Director + Sales Manager only.",
         size=9, color=C_GREY)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 6 — Award & handover
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "4. Award & Handover (Opp WON → Contract → Project)",
            "Marking an opportunity WON triggers automatic contract + project creation.")

# Center node: WON
add_box(s, Inches(5.4), Inches(1.4), Inches(2.5), Inches(0.7),
        "Opportunity → WON", fill=C_GREEN, size=12)

# Triggered records
boxes = [
    (Inches(0.6),  Inches(2.7), "Contract record\nauto-created",
     "Status = AWARDED\nValue + scope from\nfinal proposal",  C_PURPLE),
    (Inches(3.4),  Inches(2.7), "Project record\nauto-created",
     "Code · PM assigned\nGates open · RAG=Green", C_BLUE),
    (Inches(6.2),  Inches(2.7), "Milestones &\ninvoicing schedule",
     "Per contract terms\n(advance, mob, comm.)", C_AMBER),
    (Inches(9.0),  Inches(2.7), "Handover packet",
     "Notification to PM\n+ PMO Manager",          C_TEAL),
]
for x, y, label, sub, color in boxes:
    add_box(s, x, y, Inches(2.5), Inches(0.65), label, fill=color, size=10)
    add_text(s, x, y + Inches(0.7), Inches(2.5), Inches(0.7), sub,
             size=9, color=C_GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    # arrow from WON
    add_arrow(s, Inches(6.65), Inches(2.1), x + Inches(1.25), y)

# Workflow gates row
add_text(s, Inches(0.4), Inches(4.4), Inches(12.5), Inches(0.3),
         "Project gates open in sequence (PMO Manager + PM approve closure):",
         size=11, bold=True, color=C_DARK)

gates = ["G1 Kickoff", "G2 Design", "G3 Procurement", "G4 Mobilisation",
         "G5 Construction", "G6 Commissioning", "G7 Handover", "G8 Closeout"]
gn = len(gates); gw = (SLIDE_W - Inches(0.8)) / gn - Inches(0.05)
for i, g in enumerate(gates):
    x = Inches(0.4) + (gw + Inches(0.05)) * i
    add_box(s, x, Inches(4.85), gw, Inches(0.55), g, fill=C_BLUE_LT,
            text_color=C_BLUE, border=C_BLUE, size=9)
    if i < gn - 1:
        add_arrow(s, x + gw, Inches(5.125), x + gw + Inches(0.05), Inches(5.125),
                  color=C_BLUE, weight=1.25)

# Audit + notifications row
add_text(s, Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.3),
         "Triggered notifications:", size=11, bold=True, color=C_DARK)

notes = [
    ("contract_awarded",            C_GREEN_LT,  C_GREEN,
     "Director · Sales Mgr · PM · Finance"),
    ("contract_handover_ready",     C_AMBER_LT,  C_AMBER,
     "PMO Manager + PM"),
    ("contract_handover_completed", C_GREEN_LT,  C_GREEN,
     "PM + Director"),
    ("project_member_added",        C_BLUE_LT,   C_BLUE,
     "New team member"),
]
nh = Inches(0.55); ngap = Inches(0.15)
for i, (label, fill, border, sub) in enumerate(notes):
    x = Inches(0.4) + (Inches(3.05) + ngap) * i
    add_box(s, x, Inches(6.05), Inches(3.05), nh, label,
            fill=fill, text_color=border, border=border, size=9)
    add_text(s, x, Inches(6.6), Inches(3.05), Inches(0.4), sub,
             size=8, color=C_GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 7 — Project execution with gates
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "5. Project Execution — Gate, RAG & Blocker Flow",
            "Each gate has structured deliverables. Closure requires approver sign-off.")

# Left: gate state machine
add_text(s, Inches(0.4), Inches(1.2), Inches(5), Inches(0.3),
         "Single-gate state machine", size=12, bold=True, color=C_DARK)

# States
add_box(s, Inches(0.4), Inches(1.7),  Inches(2.5), Inches(0.55),
        "OPEN", fill=C_GREY, size=11)
add_box(s, Inches(0.4), Inches(2.5),  Inches(2.5), Inches(0.55),
        "IN_PROGRESS", fill=C_BLUE, size=11)
add_box(s, Inches(0.4), Inches(3.3),  Inches(2.5), Inches(0.55),
        "SUBMITTED FOR REVIEW", fill=C_AMBER, size=11)
add_diamond(s, Inches(0.4), Inches(4.1), Inches(2.5), Inches(0.9),
            "Approved?", fill=C_AMBER, size=10)
add_box(s, Inches(0.4), Inches(5.2),  Inches(2.5), Inches(0.55),
        "CLOSED", fill=C_GREEN, size=11)

# Reject loop
add_box(s, Inches(3.3), Inches(4.3),  Inches(2.0), Inches(0.55),
        "REJECTED", fill=C_RED, size=10)

add_arrow(s, Inches(1.65), Inches(2.25), Inches(1.65), Inches(2.5))
add_arrow(s, Inches(1.65), Inches(3.05), Inches(1.65), Inches(3.3))
add_arrow(s, Inches(1.65), Inches(3.85), Inches(1.65), Inches(4.1))
add_arrow(s, Inches(1.65), Inches(5.0),  Inches(1.65), Inches(5.2), color=C_GREEN)
add_text(s, Inches(1.7), Inches(5.02), Inches(1), Inches(0.2),
         "yes", size=9, bold=True, color=C_GREEN)

# reject path
add_arrow(s, Inches(2.9), Inches(4.55), Inches(3.3), Inches(4.55), color=C_RED)
add_text(s, Inches(2.95), Inches(4.32), Inches(0.4), Inches(0.2),
         "no", size=9, bold=True, color=C_RED)
# back to in_progress
add_elbow_arrow(s, Inches(3.3), Inches(4.55), Inches(3.0), Inches(2.78), color=C_RED)

# Right: project health & blocker
add_text(s, Inches(6.2), Inches(1.2), Inches(7), Inches(0.3),
         "Project-level RAG & blocker flow", size=12, bold=True, color=C_DARK)

# RAG row
add_box(s, Inches(6.2),  Inches(1.7), Inches(2),  Inches(0.55),
        "GREEN", fill=C_GREEN, size=11)
add_box(s, Inches(8.4),  Inches(1.7), Inches(2),  Inches(0.55),
        "AMBER", fill=C_AMBER, size=11)
add_box(s, Inches(10.6), Inches(1.7), Inches(2),  Inches(0.55),
        "RED", fill=C_RED, size=11)

add_text(s, Inches(6.2), Inches(2.3), Inches(6.7), Inches(0.5),
         "PM updates RAG weekly. AMBER/RED requires currentBlocker + blockerOwner + due date.",
         size=10, color=C_GREY, anchor=MSO_ANCHOR.TOP)

# Blocker assignment flow
add_box(s, Inches(6.2), Inches(3.3), Inches(2.5), Inches(0.6),
        "PM sets blocker", fill=C_BLUE, size=10)
add_box(s, Inches(9.0), Inches(3.3), Inches(2.5), Inches(0.6),
        "Assign blockerOwnerId", fill=C_BLUE, size=10)
add_box(s, Inches(11.6), Inches(3.3), Inches(1.4), Inches(0.6),
        "Notify owner", fill=C_PURPLE, size=10)
add_arrow(s, Inches(8.7), Inches(3.6), Inches(9.0), Inches(3.6))
add_arrow(s, Inches(11.5), Inches(3.6), Inches(11.6), Inches(3.6))

# Resolution
add_box(s, Inches(6.2), Inches(4.2), Inches(6.7), Inches(0.55),
        "Owner posts update → PM clears blocker → RAG recalculated",
        fill=C_GREEN_LT, text_color=C_GREEN, border=C_GREEN, size=10)
add_arrow(s, Inches(9.55), Inches(3.9), Inches(9.55), Inches(4.2))

# Deliverables panel
add_text(s, Inches(6.2), Inches(5.0), Inches(7), Inches(0.3),
         "Per-gate deliverables (examples):", size=11, bold=True, color=C_DARK)
delivs = [
    "G1: Kickoff minutes · scope confirmation · risk register",
    "G2: SLD · BoQ · DBD · permit list",
    "G3: PO list · vendor confirmations · logistics plan",
    "G4: Site readiness checklist · HSE plan · mobilisation log",
    "G5: Daily site reports · photos · punch list",
    "G6: Test sheets · SAT/FAT · commissioning checklist",
    "G7: Handover dossier · asset register · warranty docs",
    "G8: Closeout report · lessons learned · final invoice",
]
y = Inches(5.4)
for line in delivs:
    add_text(s, Inches(6.3), y, Inches(7), Inches(0.22), line,
             size=9, color=C_DARK, anchor=MSO_ANCHOR.TOP)
    y += Inches(0.22)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 8 — Procurement flow
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "6. Procurement Flow (RFQ → PO → Delivery)",
            "Triggered after Gate 2 (Design) closes. Drives Gate 3 (Procurement) closure.")

# Horizontal flow
nodes = [
    ("BoQ released\n(from Design)",     C_TEAL),
    ("Vendor list\nshortlisted",        C_TEAL),
    ("RFQ issued",                       C_BLUE),
    ("Quotes received\n+ compared",      C_BLUE),
    ("PO raised\n(approval req'd)",      C_AMBER),
    ("PO issued\nto vendor",             C_GREEN),
    ("Goods delivered\n+ GRN",           C_GREEN),
    ("Equipment serial\nrecorded",       C_PURPLE),
]
n = len(nodes)
gap = Inches(0.15)
total_w = SLIDE_W - Inches(0.8)
nw = (total_w - gap * (n - 1)) / n
nh = Inches(0.85)
ny = Inches(2.0)

for i, (label, color) in enumerate(nodes):
    x = Inches(0.4) + (nw + gap) * i
    add_box(s, x, ny, nw, nh, label, fill=color, size=10)
    if i < n - 1:
        add_arrow(s, x + nw, ny + nh/2,
                  x + nw + gap, ny + nh/2)

# Decision: PO approval
add_diamond(s, Inches(5.5), Inches(3.6), Inches(2.4), Inches(1.0),
            "Value > approval\nthreshold?",
            fill=C_AMBER, size=10)
# arrow up from RFQ→PO chain into diamond (visually annotation)
add_arrow(s, Inches(6.7), Inches(2.85), Inches(6.7), Inches(3.6), color=C_AMBER, weight=1.5)

# Yes → Director approves
add_box(s, Inches(8.3), Inches(3.85), Inches(3.0), Inches(0.55),
        "Director approves", fill=C_PURPLE, size=10)
add_arrow(s, Inches(7.9), Inches(4.1), Inches(8.3), Inches(4.1), color=C_PURPLE)
add_text(s, Inches(7.92), Inches(3.85), Inches(0.4), Inches(0.25),
         "yes", size=9, bold=True, color=C_PURPLE)

# No → straight to issue
add_box(s, Inches(2.3), Inches(3.85), Inches(2.7), Inches(0.55),
        "PM signs off", fill=C_BLUE, size=10)
add_arrow(s, Inches(5.5), Inches(4.1), Inches(5.0), Inches(4.1), color=C_BLUE)
add_text(s, Inches(5.05), Inches(3.85), Inches(0.4), Inches(0.25),
         "no", size=9, bold=True, color=C_BLUE)

# Status notifications
add_text(s, Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.3),
         "PO status changes notify: project PM · project members · all PROCUREMENT-role users",
         size=11, bold=True, color=C_DARK)

# Status pills row
states = [
    ("DRAFT",       C_GREY_LT,  C_GREY),
    ("SUBMITTED",   C_BLUE_LT,  C_BLUE),
    ("APPROVED",    C_GREEN_LT, C_GREEN),
    ("ISSUED",      C_PURPLE_LT,C_PURPLE),
    ("PARTIAL DELIVERED", C_AMBER_LT, C_AMBER),
    ("DELIVERED",   C_GREEN_LT, C_GREEN),
    ("CLOSED",      C_GREY_LT,  C_GREY),
    ("CANCELLED",   C_RED_LT,   C_RED),
]
sw = Inches(1.5); sgap = Inches(0.08)
for i, (label, fill, border) in enumerate(states):
    x = Inches(0.4) + (sw + sgap) * i
    add_box(s, x, Inches(5.45), sw, Inches(0.5), label,
            fill=fill, text_color=border, border=border, size=9)

# Asset registration foot
add_box(s, Inches(0.4), Inches(6.3), SLIDE_W - Inches(0.8), Inches(0.5),
        "On final delivery: serials captured per item → ready for asset register at commissioning",
        fill=C_TEAL_LT, text_color=C_TEAL, border=C_TEAL, size=11)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 9 — Commissioning → Asset Register → O&M
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "7. Commissioning → Asset Register → O&M",
            "Where the project becomes an operating asset under maintenance.")

# Three column flow
col_titles = [
    ("Commissioning", C_AMBER),
    ("Handover & Asset Register", C_TEAL),
    ("O&M Operating", C_GREEN),
]
col_w = Inches(4.1); col_gap = Inches(0.15)
for i, (title, color) in enumerate(col_titles):
    x = Inches(0.4) + (col_w + col_gap) * i
    add_box(s, x, Inches(1.2), col_w, Inches(0.55), title, fill=color, size=14)

# Commissioning steps
commis_steps = [
    "Pre-commissioning checks",
    "Functional test sheets",
    "SAT / FAT execution",
    "Punch list cleared",
    "Gate 6 closed (PMO approval)",
]
y = Inches(1.95)
for step in commis_steps:
    add_box(s, Inches(0.4), y, col_w, Inches(0.5), step,
            fill=C_AMBER_LT, text_color=C_AMBER, border=C_AMBER, size=10)
    y += Inches(0.62)

# Handover & asset
ho_steps = [
    "Asset register populated\n(serial · warranty start · location)",
    "Customer handover dossier",
    "Customer sign-off captured",
    "Gate 7 closed",
    "Notify Finance for final invoice",
]
y = Inches(1.95); x = Inches(0.4) + (col_w + col_gap)
for step in ho_steps:
    add_box(s, x, y, col_w, Inches(0.5), step,
            fill=C_TEAL_LT, text_color=C_TEAL, border=C_TEAL, size=10)
    y += Inches(0.62)

# O&M
om_steps = [
    "Maintenance plan auto-created\n(per asset class)",
    "Service tickets raised\nas issues arise",
    "Scheduled maintenance executed",
    "Warranty claims filed if needed",
    "Asset history & uptime tracked",
]
y = Inches(1.95); x = Inches(0.4) + (col_w + col_gap) * 2
for step in om_steps:
    add_box(s, x, y, col_w, Inches(0.5), step,
            fill=C_GREEN_LT, text_color=C_GREEN, border=C_GREEN, size=10)
    y += Inches(0.62)

# Inter-column arrows (between last commissioning step and first handover, etc.)
add_arrow(s, Inches(0.4) + col_w, Inches(2.2), Inches(0.4) + col_w + col_gap, Inches(2.2),
          color=C_GREY, weight=2)
add_arrow(s, Inches(0.4) + col_w*2 + col_gap, Inches(2.2),
          Inches(0.4) + col_w*2 + col_gap*2, Inches(2.2),
          color=C_GREY, weight=2)

# Bottom: who does what
add_text(s, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.3),
         "Owners & data flow:", size=11, bold=True, color=C_DARK)

owners = [
    ("Commissioning Engineer",   "Runs all test sheets · marks pass/fail · uploads evidence"),
    ("PM + PMO Manager",          "Approve gate closure · sign off handover packet"),
    ("Site Supervisor",           "Captures serials & install location at GRN time"),
    ("O&M Engineer",              "Owns asset post-handover · all tickets routed here"),
    ("Finance/Admin",             "Triggers final milestone invoice on Gate 7 closure"),
]
y = Inches(6.2)
for role, what in owners:
    add_text(s, Inches(0.5), y, Inches(2.3), Inches(0.25),
             "• " + role, size=10, bold=True, color=C_DARK, anchor=MSO_ANCHOR.TOP)
    add_text(s, Inches(2.9), y, Inches(10), Inches(0.25),
             what, size=10, color=C_GREY, anchor=MSO_ANCHOR.TOP)
    y += Inches(0.22)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 10 — Cross-cutting: notifications, audit, role visibility
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_title(s, "8. Cross-Cutting Layers — Notifications · Audit · Role Visibility",
            "Same data, different views. The system enforces who sees what at every stage.")

# Row 1: notification triggers
add_text(s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.3),
         "Notification triggers (in-app bell + daily digest):",
         size=12, bold=True, color=C_DARK)

triggers = [
    ("Gate submitted/approved/rejected", C_AMBER, "Gate approver + PM"),
    ("Proposal submitted/decision",      C_BLUE,  "Sales mgr + Director"),
    ("Opp WON / LOST",                   C_GREEN, "Sales team + Director"),
    ("Deliverable approved",             C_GREEN, "PM + project members"),
    ("Blocker assigned",                 C_RED,   "Blocker owner only"),
    ("PO status changed",                C_AMBER, "PM + PROCUREMENT users"),
    ("Document approved/rejected",       C_TEAL,  "Uploader + doc owner"),
    ("Project member added",             C_BLUE,  "New member + PM"),
    ("Overdue next action",              C_RED,   "Action + opp owner"),
    ("Proposal awaiting follow-up (3d)", C_AMBER, "Opp owner"),
    ("Stale opportunity (14d)",          C_AMBER, "Opp owner"),
    ("Daily digest",                     C_PURPLE,"Each user (opt-in)"),
]
tw = Inches(3.1); tgap = Inches(0.12)
cols = 4
for idx, (trig, color, who) in enumerate(triggers):
    row = idx // cols
    col = idx % cols
    x = Inches(0.4) + (tw + tgap) * col
    y = Inches(1.6) + Inches(0.85) * row
    add_box(s, x, y, tw, Inches(0.45), trig, fill=color, size=9)
    add_text(s, x, y + Inches(0.5), tw, Inches(0.3), who,
             size=8, color=C_GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# Row 2: role visibility matrix (compact)
add_text(s, Inches(0.4), Inches(4.85), Inches(12.5), Inches(0.3),
         "Role visibility (sensitive fields):",
         size=12, bold=True, color=C_DARK)

# header row
roles_hdr = ["", "Director", "Sales Mgr", "Sales Eng", "PMO Mgr", "PM", "Design", "Procure", "Site Sup", "O&M Eng", "Finance"]
fields = [
    ("Margin",          ["✓","✓","",  "",  "",  "",  "",  "",  "",  "✓"]),
    ("CAPEX value",     ["✓","✓","",  "",  "",  "",  "",  "",  "",  "✓"]),
    ("Pipeline value",  ["✓","✓","",  "",  "",  "",  "",  "",  "",  ""]),
    ("Opportunity $",   ["✓","✓","✓", "",  "",  "",  "",  "",  "",  "✓"]),
    ("Project gate",    ["✓","",   "",  "✓", "✓", "✓", "✓", "✓", "",  ""]),
    ("Asset register",  ["✓","",   "",  "✓", "✓", "",  "",  "✓", "✓", ""]),
    ("Audit log",       ["✓","",   "",  "✓", "",  "",  "",  "",  "",  ""]),
]

table_y = Inches(5.25)
col_w_tbl = Inches(1.05); col_w_first = Inches(1.6)
row_h = Inches(0.32)

# header
for i, h in enumerate(roles_hdr):
    x = Inches(0.4) + (col_w_first if i > 0 else Inches(0)) + col_w_tbl * (i - 1) if i > 0 else Inches(0.4)
    if i == 0:
        cw = col_w_first
    else:
        cw = col_w_tbl
    cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_y, cw, row_h)
    cell.fill.solid(); cell.fill.fore_color.rgb = C_BLUE
    cell.line.color.rgb = C_WHITE
    cell.shadow.inherit = False
    tf = cell.text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = h
    r.font.size = Pt(8); r.font.bold = True; r.font.color.rgb = C_WHITE
    r.font.name = "Calibri"

# rows
for ri, (fname, marks) in enumerate(fields):
    yr = table_y + row_h * (ri + 1)
    # field cell
    cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), yr, col_w_first, row_h)
    cell.fill.solid(); cell.fill.fore_color.rgb = C_GREY_LT
    cell.line.color.rgb = C_WHITE
    cell.shadow.inherit = False
    tf = cell.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = " " + fname
    r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = C_DARK
    r.font.name = "Calibri"
    for ci, mark in enumerate(marks):
        x = Inches(0.4) + col_w_first + col_w_tbl * ci
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, yr, col_w_tbl, row_h)
        if mark == "✓":
            cell.fill.solid(); cell.fill.fore_color.rgb = C_GREEN_LT
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = C_WHITE
        cell.line.color.rgb = C_GREY_LT
        cell.shadow.inherit = False
        tf = cell.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = mark
        r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = C_GREEN if mark == "✓" else C_GREY_LT
        r.font.name = "Calibri"

# Save
prs.save(OUT)
print(f"Wrote {OUT}")
print(f"Size: {os.path.getsize(OUT)} bytes")
